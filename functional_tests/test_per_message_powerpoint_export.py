#!/usr/bin/env python3
# test_per_message_powerpoint_export.py
"""
Functional test for per-message PowerPoint export.
Version: 0.241.105
Implemented in: 0.241.105

This test ensures the message export flow exposes a PowerPoint route,
uses the frontend PowerPoint action hook, prefers the message model
deployment for AI slide planning, and produces a valid .pptx deck with
appendix slides for visuals, tables, code, and references.
"""

import ast
import base64
import io
import json
import os
import re
import traceback
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Dict, List, Optional, Tuple

import markdown2
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt


REPO_ROOT = Path(__file__).resolve().parents[1]
ROUTE_FILE = REPO_ROOT / "application" / "single_app" / "route_backend_conversation_export.py"
FRONTEND_FILE = REPO_ROOT / "application" / "single_app" / "static" / "js" / "chat" / "chat-message-export.js"
MENU_FILE = REPO_ROOT / "application" / "single_app" / "static" / "js" / "chat" / "chat-messages.js"


def _normalize_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict):
                if item.get('type') == 'text':
                    parts.append(item.get('text', ''))
                elif item.get('type') == 'image_url':
                    parts.append('[Image]')
                else:
                    parts.append(str(item))
            else:
                parts.append(str(item))
        return '\n'.join(parts)
    if isinstance(content, dict):
        if content.get('type') == 'text':
            return content.get('text', '')
        return str(content)
    return str(content) if content else ''


def _decode_base64_image_data_uri(data_uri: Optional[str]) -> Optional[bytes]:
    if not data_uri or not isinstance(data_uri, str):
        return None

    match = re.match(r'^data:image\/[a-zA-Z0-9.+-]+;base64,(.+)$', data_uri.strip())
    if not match:
        return None

    try:
        return base64.b64decode(match.group(1))
    except Exception:
        return None


def _build_test_image_data_uri() -> str:
    image = Image.new('RGB', (24, 16), color=(37, 99, 235))
    buffer = io.BytesIO()
    image.save(buffer, format='PNG')
    encoded = base64.b64encode(buffer.getvalue()).decode('ascii')
    return f'data:image/png;base64,{encoded}'


def _collect_slide_titles(presentation: Presentation) -> List[str]:
    titles = []
    for slide in presentation.slides:
        title_shape = slide.shapes.title
        if title_shape and title_shape.text:
            titles.append(title_shape.text.strip())
    return titles


def _load_powerpoint_helpers():
    helper_names = {
        '_message_to_pptx_bytes',
        '_sanitize_powerpoint_source_content',
        '_build_message_powerpoint_plan',
        '_build_fallback_powerpoint_plan',
        '_extract_message_powerpoint_model',
        '_normalize_powerpoint_model_candidate',
        '_generate_powerpoint_slide_plan_with_model',
        '_extract_json_object',
        '_sanitize_powerpoint_plan',
        '_extract_powerpoint_sections',
        '_extract_powerpoint_bullets',
        '_sentence_bullets',
        '_looks_like_markdown_table_row',
        '_looks_like_markdown_table_divider',
        '_markdown_to_plain_text',
        '_derive_powerpoint_title',
        '_build_powerpoint_subtitle',
        '_clean_slide_text',
        '_extract_powerpoint_appendix_assets',
        '_extract_powerpoint_images',
        '_extract_powerpoint_tables',
        '_extract_powerpoint_code_blocks',
        '_add_powerpoint_title_slide',
        '_add_powerpoint_content_slide',
        '_append_powerpoint_appendix_slides',
        '_add_powerpoint_image_slide',
        '_fit_powerpoint_image',
        '_add_powerpoint_table_slide',
        '_add_powerpoint_code_slide',
        '_apply_powerpoint_background',
        '_chunk_items',
    }

    source = ROUTE_FILE.read_text(encoding='utf-8')
    tree = ast.parse(source)
    selected_nodes = [
        node for node in tree.body
        if isinstance(node, ast.FunctionDef) and node.name in helper_names
    ]

    loaded_names = {node.name for node in selected_nodes}
    missing_names = helper_names - loaded_names
    assert not missing_names, f"Missing PowerPoint helpers in route file: {sorted(missing_names)}"

    requested_models: List[str] = []

    class _FakeCompletions:
        def create(self, model, messages):
            slide_plan = {
                'presentation_title': 'Quarterly Review',
                'presentation_subtitle': 'Assistant | Generated from chat message',
                'slides': [
                    {
                        'title': 'Overview',
                        'bullets': [
                            'Revenue grew 18 percent.',
                            'Support backlog dropped 12 percent.',
                        ],
                    },
                    {
                        'title': 'Presenter Notes',
                        'bullets': [
                            'Use the visual, table, and code appendix as backup.',
                        ],
                    },
                ],
            }
            return SimpleNamespace(
                choices=[
                    SimpleNamespace(
                        message=SimpleNamespace(content=json.dumps(slide_plan))
                    )
                ]
            )

    class _FakeClient:
        def __init__(self):
            self.chat = SimpleNamespace(completions=_FakeCompletions())

    def _fake_initialize_gpt_client(settings, requested_model=''):
        requested_models.append(requested_model or '')
        return _FakeClient(), requested_model or 'fallback-model'

    module = ast.Module(body=selected_nodes, type_ignores=[])
    ast.fix_missing_locations(module)

    namespace = {
        'Any': Any,
        'Dict': Dict,
        'List': List,
        'Optional': Optional,
        'Tuple': Tuple,
        'io': io,
        'json': json,
        're': re,
        'markdown2': markdown2,
        'BeautifulSoup': BeautifulSoup,
        'NavigableString': NavigableString,
        'Tag': Tag,
        'Image': Image,
        'Presentation': Presentation,
        'RGBColor': RGBColor,
        'MSO_AUTO_SHAPE_TYPE': MSO_AUTO_SHAPE_TYPE,
        'PP_ALIGN': PP_ALIGN,
        'PptxInches': PptxInches,
        'PptxPt': PptxPt,
        'DOCX_MARKDOWN_EXTRAS': ['fenced-code-blocks', 'tables', 'break-on-newline', 'cuddled-lists', 'strike'],
        'POWERPOINT_PLAN_SOURCE_CHAR_LIMIT': 24000,
        'POWERPOINT_MAX_SLIDES': 7,
        'POWERPOINT_MAX_BULLETS_PER_SLIDE': 5,
        'POWERPOINT_MAX_APPENDIX_IMAGES': 4,
        'POWERPOINT_MAX_APPENDIX_TABLES': 3,
        'POWERPOINT_MAX_APPENDIX_CODE_BLOCKS': 2,
        'POWERPOINT_MAX_TABLE_ROWS': 8,
        'POWERPOINT_MAX_TABLE_COLS': 5,
        'POWERPOINT_TITLE_BG': RGBColor(22, 37, 66),
        'POWERPOINT_ACCENT': RGBColor(37, 99, 235),
        'POWERPOINT_BG': RGBColor(248, 250, 252),
        'POWERPOINT_PANEL': RGBColor(255, 255, 255),
        'POWERPOINT_TEXT': RGBColor(31, 41, 55),
        'POWERPOINT_MUTED': RGBColor(100, 116, 139),
        'POWERPOINT_TITLE_TEXT': RGBColor(255, 255, 255),
        'POWERPOINT_DATA_URI_PATTERN': re.compile(
            r"data:image\/[a-zA-Z0-9.+-]+;base64,[^\"'\s)]+",
            re.IGNORECASE,
        ),
        '_normalize_content': _normalize_content,
        '_role_to_label': lambda role: {
            'assistant': 'Assistant',
            'user': 'User',
            'system': 'System',
        }.get(role, str(role).capitalize() or 'Message'),
        '_build_message_citation_labels': lambda message: [
            citation.get('title') or citation.get('label') or citation.get('url') or str(citation)
            for citation in message.get('citations', [])
        ],
        'replace_inline_chart_blocks_with_export_html': lambda content: content,
        'decode_base64_image_data_uri': _decode_base64_image_data_uri,
        '_initialize_gpt_client': _fake_initialize_gpt_client,
        'debug_print': lambda *args, **kwargs: None,
        'log_event': lambda *args, **kwargs: None,
    }

    exec(compile(module, str(ROUTE_FILE), 'exec'), namespace)
    return namespace, requested_models


def test_export_powerpoint_route_definition_present() -> bool:
    """Route regression: the backend must define POST /api/message/export-powerpoint."""
    print("Testing backend route definition for PowerPoint export...")

    source = ROUTE_FILE.read_text(encoding='utf-8')
    tree = ast.parse(source)
    register_func = next(
        (
            node for node in tree.body
            if isinstance(node, ast.FunctionDef)
            and node.name == 'register_route_backend_conversation_export'
        ),
        None,
    )

    assert register_func is not None, 'register_route_backend_conversation_export should exist'

    route_found = False
    for node in register_func.body:
        if not isinstance(node, ast.FunctionDef):
            continue

        for decorator in node.decorator_list:
            if not isinstance(decorator, ast.Call):
                continue
            if not isinstance(decorator.func, ast.Attribute) or decorator.func.attr != 'route':
                continue
            if not decorator.args:
                continue

            route_arg = decorator.args[0]
            if not isinstance(route_arg, ast.Constant) or route_arg.value != '/api/message/export-powerpoint':
                continue

            methods_kw = next((keyword for keyword in decorator.keywords if keyword.arg == 'methods'), None)
            assert methods_kw is not None, 'PowerPoint export route should declare allowed methods'
            methods = [
                item.value for item in methods_kw.value.elts
                if isinstance(item, ast.Constant)
            ]
            assert 'POST' in methods, f'Expected POST method, found {methods}'
            assert node.name == 'api_export_message_powerpoint', f'Unexpected route handler name: {node.name}'
            route_found = True
            break

        if route_found:
            break

    assert route_found, 'Expected POST /api/message/export-powerpoint to be defined'
    print("PASS: backend route definition present")
    return True


def test_powerpoint_frontend_hooks_present() -> bool:
    """Frontend regression: menu and fetch path should expose PowerPoint export."""
    print("Testing frontend PowerPoint export hooks...")

    frontend_source = FRONTEND_FILE.read_text(encoding='utf-8')
    menu_source = MENU_FILE.read_text(encoding='utf-8')

    assert "fetch('/api/message/export-powerpoint'" in frontend_source, 'Expected frontend fetch for the PowerPoint export endpoint'
    assert 'exportMessageAsPowerPoint' in frontend_source, 'Expected frontend PowerPoint export helper'
    assert 'dropdown-export-ppt-btn' in menu_source, 'Expected chat message menu PowerPoint action'

    print("PASS: frontend PowerPoint hooks present")
    return True


def test_powerpoint_export_prefers_message_model_and_renders_appendix() -> bool:
    """PowerPoint export should use the message model hint and render appendix slides."""
    print("Testing PowerPoint slide generation...")

    helpers, requested_models = _load_powerpoint_helpers()
    image_data_uri = _build_test_image_data_uri()
    message = {
        'role': 'assistant',
        'timestamp': '2026-05-04T12:00:00Z',
        'model_deployment_name': 'gpt-4o-mini',
        'content': '\n'.join([
            '# Quarterly Review',
            '',
            '- Revenue grew 18 percent.',
            '- Support backlog dropped 12 percent.',
            '',
            '| Metric | Value |',
            '| --- | --- |',
            '| Revenue | +18% |',
            '| Backlog | -12% |',
            '',
            '```python',
            'print("hello slides")',
            '```',
            '',
            f'<div class="export-inline-chart"><img src="{image_data_uri}" alt="Trend chart" /><div class="export-inline-chart-caption">Trend chart</div></div>',
        ]),
        'citations': [
            {'title': 'Quarterly workbook'},
            {'title': 'Operations dashboard'},
        ],
    }

    pptx_bytes = helpers['_message_to_pptx_bytes'](message, {'gpt_model': {'selected': [{'deploymentName': 'fallback-model'}]}})

    assert requested_models == ['gpt-4o-mini'], f'Expected the message deployment to be reused, found {requested_models}'
    assert pptx_bytes[:2] == b'PK', 'PowerPoint export should return a zipped OOXML payload'

    presentation = Presentation(io.BytesIO(pptx_bytes))
    slide_titles = _collect_slide_titles(presentation)

    assert len(presentation.slides) >= 6, f'Expected multiple slides including appendix content, found {len(presentation.slides)}'
    assert slide_titles[0] == 'Quarterly Review', f'Unexpected title slide text: {slide_titles[0]}'
    assert 'Overview' in slide_titles, f'Missing AI outline slide in {slide_titles}'
    assert 'Presenter Notes' in slide_titles, f'Missing second AI outline slide in {slide_titles}'
    assert 'Visual 1' in slide_titles, f'Missing visual appendix slide in {slide_titles}'
    assert 'Table 1' in slide_titles, f'Missing table appendix slide in {slide_titles}'
    assert 'Code Example 1' in slide_titles, f'Missing code appendix slide in {slide_titles}'
    assert 'References' in slide_titles, f'Missing references slide in {slide_titles}'

    print("PASS: PowerPoint export renders slide deck with appendix content")
    return True


if __name__ == '__main__':
    tests = [
        test_export_powerpoint_route_definition_present,
        test_powerpoint_frontend_hooks_present,
        test_powerpoint_export_prefers_message_model_and_renders_appendix,
    ]

    results = []
    for test in tests:
        print(f"\nRunning {test.__name__}...")
        try:
            results.append(bool(test()))
        except Exception as exc:
            print(f"FAIL: {test.__name__}: {exc}")
            traceback.print_exc()
            results.append(False)

    passed = sum(1 for result in results if result)
    print(f"\nResults: {passed}/{len(results)} tests passed")
    raise SystemExit(0 if all(results) else 1)